# ============================================================
# DOCUMENT SERVICE
# ============================================================
#
# Purpose:
#   Extract readable text from uploaded learning documents
#   so that the AI Quiz Generation Service can use the content.
#
# Supported formats:
#   - PDF
#   - PPTX
#   - TXT
#
# Flow:
#
#   Uploaded File
#        ↓
#   validate file type
#        ↓
#   extract text
#        ↓
#   clean text
#        ↓
#   validate extracted content
#        ↓
#   return text
#
# This service does NOT:
#   - generate questions
#   - call Gemini
#   - create assessments
#   - modify database records
#
# Those responsibilities belong to later services.
# ============================================================

from __future__ import annotations

import io
import re
from pathlib import Path
from typing import BinaryIO


# ============================================================
# CONSTANTS
# ============================================================

SUPPORTED_EXTENSIONS = {
    ".pdf",
    ".pptx",
    ".txt",
}

MAX_DOCUMENT_SIZE_BYTES = 20 * 1024 * 1024  # 20 MB

MIN_EXTRACTED_TEXT_LENGTH = 50


# ============================================================
# CUSTOM EXCEPTIONS
# ============================================================


class DocumentServiceError(Exception):
    """
    Base exception for document processing errors.
    """

    pass


class UnsupportedDocumentTypeError(DocumentServiceError):
    """
    Raised when the uploaded document format is unsupported.
    """

    pass


class DocumentTooLargeError(DocumentServiceError):
    """
    Raised when the uploaded document exceeds the allowed size.
    """

    pass


class DocumentExtractionError(DocumentServiceError):
    """
    Raised when text extraction fails.
    """

    pass


class EmptyDocumentError(DocumentServiceError):
    """
    Raised when no useful text can be extracted.
    """

    pass


# ============================================================
# FILE TYPE VALIDATION
# ============================================================


def get_file_extension(filename: str) -> str:
    """
    Return the normalized file extension.

    Example:
        "training_material.PDF" → ".pdf"
    """

    if not filename:
        raise UnsupportedDocumentTypeError(
            "Document filename is required."
        )

    extension = Path(filename).suffix.lower()

    if extension not in SUPPORTED_EXTENSIONS:
        supported = ", ".join(sorted(SUPPORTED_EXTENSIONS))

        raise UnsupportedDocumentTypeError(
            f"Unsupported document type '{extension or 'unknown'}'. "
            f"Supported formats: {supported}"
        )

    return extension


# ============================================================
# SIZE VALIDATION
# ============================================================


def validate_document_size(content: bytes) -> None:
    """
    Validate uploaded document size.

    The quiz generator does not need extremely large files.
    Limiting the size also prevents unnecessarily large AI
    prompts and protects backend resources.
    """

    if not content:
        raise DocumentExtractionError(
            "Uploaded document is empty."
        )

    if len(content) > MAX_DOCUMENT_SIZE_BYTES:
        size_mb = len(content) / (1024 * 1024)

        raise DocumentTooLargeError(
            f"Document size is {size_mb:.2f} MB. "
            f"Maximum allowed size is "
            f"{MAX_DOCUMENT_SIZE_BYTES / (1024 * 1024):.0f} MB."
        )


# ============================================================
# TEXT CLEANING
# ============================================================


def clean_extracted_text(text: str) -> str:
    """
    Normalize extracted document text.

    Cleaning includes:
        - removing null characters
        - normalizing line endings
        - removing excessive whitespace
        - removing excessive blank lines
        - preserving paragraph structure
    """

    if not text:
        return ""

    # Remove null/control characters that can appear in
    # extracted PDF/PPT content.
    text = text.replace("\x00", " ")

    # Normalize line endings.
    text = text.replace("\r\n", "\n")
    text = text.replace("\r", "\n")

    # Normalize tabs.
    text = text.replace("\t", " ")

    # Remove spaces at the end of lines.
    text = re.sub(r"[ \t]+\n", "\n", text)

    # Collapse repeated spaces while preserving newlines.
    text = re.sub(r"[ \t]{2,}", " ", text)

    # Collapse excessive blank lines.
    text = re.sub(r"\n{3,}", "\n\n", text)

    return text.strip()


# ============================================================
# TXT EXTRACTION
# ============================================================


def extract_text_from_txt(content: bytes) -> str:
    """
    Extract text from a plain-text document.

    UTF-8 is attempted first.
    A fallback decoding is used for documents that contain
    common Windows/legacy characters.
    """

    try:
        return content.decode("utf-8")

    except UnicodeDecodeError:
        try:
            return content.decode("utf-8-sig")

        except UnicodeDecodeError:
            try:
                return content.decode("cp1252")

            except UnicodeDecodeError as exc:
                raise DocumentExtractionError(
                    "Unable to decode the TXT document."
                ) from exc


# ============================================================
# PDF EXTRACTION
# ============================================================


def extract_text_from_pdf(content: bytes) -> str:
    """
    Extract text from a PDF document.

    Uses pypdf.

    Note:
        This extracts machine-readable text.
        Scanned/image-only PDFs require OCR and are intentionally
        not handled by this first MVP implementation.
    """

    try:
        from pypdf import PdfReader

    except ImportError as exc:
        raise DocumentExtractionError(
            "PDF support requires the 'pypdf' package. "
            "Install it using: pip install pypdf"
        ) from exc

    try:
        reader = PdfReader(io.BytesIO(content))

        extracted_pages: list[str] = []

        for page in reader.pages:
            try:
                page_text = page.extract_text() or ""

            except Exception:
                # One problematic page should not necessarily
                # destroy extraction from the entire document.
                page_text = ""

            if page_text.strip():
                extracted_pages.append(page_text)

        return "\n\n".join(extracted_pages)

    except Exception as exc:
        raise DocumentExtractionError(
            "Failed to extract text from the PDF document."
        ) from exc


# ============================================================
# PPTX EXTRACTION
# ============================================================


def extract_text_from_pptx(content: bytes) -> str:
    """
    Extract readable text from a PowerPoint PPTX document.

    Text is collected from:
        - slide shapes
        - text frames
        - tables
    """

    try:
        from pptx import Presentation

    except ImportError as exc:
        raise DocumentExtractionError(
            "PPTX support requires the 'python-pptx' package. "
            "Install it using: pip install python-pptx"
        ) from exc

    try:
        presentation = Presentation(io.BytesIO(content))

        slides_text: list[str] = []

        for slide_number, slide in enumerate(
            presentation.slides,
            start=1,
        ):
            slide_parts: list[str] = []

            for shape in slide.shapes:

                # ------------------------------------------------
                # Normal text boxes / placeholders
                # ------------------------------------------------

                if getattr(shape, "has_text_frame", False):
                    text = shape.text or ""

                    if text.strip():
                        slide_parts.append(text)

                # ------------------------------------------------
                # Tables
                # ------------------------------------------------

                if getattr(shape, "has_table", False):
                    table = shape.table

                    for row in table.rows:
                        row_values: list[str] = []

                        for cell in row.cells:
                            cell_text = cell.text or ""

                            if cell_text.strip():
                                row_values.append(cell_text.strip())

                        if row_values:
                            slide_parts.append(
                                " | ".join(row_values)
                            )

            if slide_parts:
                slides_text.append(
                    f"Slide {slide_number}\n"
                    + "\n".join(slide_parts)
                )

        return "\n\n".join(slides_text)

    except Exception as exc:
        raise DocumentExtractionError(
            "Failed to extract text from the PPTX document."
        ) from exc


# ============================================================
# DOCUMENT TEXT EXTRACTION
# ============================================================


def extract_document_text(
    content: bytes,
    filename: str,
) -> str:
    """
    Extract and clean text from a supported document.

    Args:
        content:
            Raw uploaded file bytes.

        filename:
            Original filename including extension.

    Returns:
        Cleaned document text.

    Raises:
        UnsupportedDocumentTypeError
        DocumentTooLargeError
        DocumentExtractionError
        EmptyDocumentError
    """

    extension = get_file_extension(filename)

    validate_document_size(content)

    try:

        if extension == ".txt":
            text = extract_text_from_txt(content)

        elif extension == ".pdf":
            text = extract_text_from_pdf(content)

        elif extension == ".pptx":
            text = extract_text_from_pptx(content)

        else:
            # This should never happen because extension validation
            # already occurs above.
            raise UnsupportedDocumentTypeError(
                f"Unsupported document extension: {extension}"
            )

    except DocumentServiceError:
        raise

    except Exception as exc:
        raise DocumentExtractionError(
            "An unexpected error occurred while processing "
            "the document."
        ) from exc

    cleaned_text = clean_extracted_text(text)

    if len(cleaned_text) < MIN_EXTRACTED_TEXT_LENGTH:
        raise EmptyDocumentError(
            "The document does not contain enough readable text "
            "to generate a meaningful quiz."
        )

    return cleaned_text


# ============================================================
# FILE OBJECT SUPPORT
# ============================================================


def extract_document_text_from_file(
    file: BinaryIO,
    filename: str,
) -> str:
    """
    Convenience helper for FastAPI UploadFile/file objects.

    The function reads the file into memory and passes it to
    the main extraction pipeline.
    """

    try:
        content = file.read()

    except Exception as exc:
        raise DocumentExtractionError(
            "Unable to read the uploaded document."
        ) from exc

    if not isinstance(content, bytes):
        raise DocumentExtractionError(
            "Uploaded document content must be bytes."
        )

    return extract_document_text(
        content=content,
        filename=filename,
    )


# ============================================================
# DOCUMENT METADATA
# ============================================================


def get_document_info(
    filename: str,
    content: bytes,
) -> dict:
    """
    Return basic document metadata.

    This is useful for the quiz-generation endpoint and logging.
    """

    extension = get_file_extension(filename)

    return {
        "filename": filename,
        "extension": extension,
        "size_bytes": len(content),
        "size_mb": round(
            len(content) / (1024 * 1024),
            2,
        ),
        "supported": True,
    }